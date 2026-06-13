#!/usr/bin/env python3
"""Generate a readable and objective PPTX deck for Lab03 S03."""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent
ASSETS_DIR = SCRIPT_DIR / "assets"
DEFAULT_OUTPUT = SCRIPT_DIR / "lab03s03-apresentacao.pptx"

BG_COLOR = RGBColor(7, 22, 49)
SURFACE_COLOR = RGBColor(18, 45, 84)
CARD_COLOR = RGBColor(23, 56, 101)
ACCENT_COLOR = RGBColor(27, 183, 129)
TEXT_PRIMARY = RGBColor(244, 248, 255)
TEXT_SECONDARY = RGBColor(184, 203, 231)
MUTED_LINE = RGBColor(63, 94, 140)

IMG_BOXES = {
    "landscape": (7.75, 1.65, 4.85, 3.20),
    "portrait": (8.05, 1.55, 4.45, 4.95),
    "diagram": (7.85, 1.55, 4.70, 4.95),
}


@dataclass(frozen=True)
class CardBlock:
    title: str
    items: list[str]


@dataclass(frozen=True)
class SplitSlide:
    title: str
    bullets: list[str]
    image: str | None = None
    image_mode: str = "landscape"
    footer: str = "Sistema de Moeda Estudantil | Lab03"


def disable_spellcheck(paragraph) -> None:
    for run in paragraph.runs:
        rpr = run._r.get_or_add_rPr()
        if rpr.find(qn("a:noProof")) is None:
            rpr.append(OxmlElement("a:noProof"))
        rpr.set("lang", "pt-BR")


def style_paragraph(
    paragraph,
    *,
    size: int,
    color: RGBColor,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 1.15,
    space_after: int = 6,
) -> None:
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    paragraph.space_after = Pt(space_after)
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    paragraph.font.name = "Aptos"
    disable_spellcheck(paragraph)


def set_background(slide, prs: Presentation) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.18))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ACCENT_COLOR
    top_bar.line.fill.background()


def add_footer(slide, prs: Presentation, text: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(0.45),
        prs.slide_height - Inches(0.38),
        prs.slide_width - Inches(0.9),
        Inches(0.22),
    )
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = text
    style_paragraph(p, size=10, color=TEXT_SECONDARY, align=PP_ALIGN.RIGHT, space_after=0)


def add_title_block(slide, title: str, subtitle: str) -> None:
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.6),
        Inches(0.62),
        Inches(3.1),
        Inches(0.42),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = SURFACE_COLOR
    tag.line.color.rgb = ACCENT_COLOR

    tag_tf = tag.text_frame
    tag_tf.clear()
    p_tag = tag_tf.paragraphs[0]
    p_tag.text = "LAB03 | APRESENTACAO S03"
    style_paragraph(p_tag, size=12, color=TEXT_SECONDARY, bold=True, align=PP_ALIGN.CENTER, space_after=0)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.25), Inches(6.7), Inches(1.2))
    title_tf = title_box.text_frame
    title_tf.clear()
    p_title = title_tf.paragraphs[0]
    p_title.text = title
    style_paragraph(p_title, size=40, color=TEXT_PRIMARY, bold=True, space_after=0)

    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.38), Inches(6.8), Inches(1.15))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    p_sub = subtitle_tf.paragraphs[0]
    p_sub.text = subtitle
    style_paragraph(p_sub, size=19, color=TEXT_SECONDARY, line_spacing=1.2, space_after=0)


def add_image_panel(slide, image_path: Path, box: tuple[float, float, float, float]) -> None:
    left_in, top_in, width_in, height_in = box
    left = Inches(left_in)
    top = Inches(top_in)
    width = Inches(width_in)
    height = Inches(height_in)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE_COLOR
    panel.line.color.rgb = MUTED_LINE
    panel.line.width = Pt(1.2)

    if not image_path.exists():
        tf = panel.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"Imagem nao encontrada: {image_path.name}"
        style_paragraph(p, size=12, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        return

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    box_w = float(width)
    box_h = float(height)
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    margin = float(Inches(0.12))

    if img_ratio > box_ratio:
        render_w = box_w - 2 * margin
        render_h = render_w / img_ratio
        render_x = float(left) + margin
        render_y = float(top) + (box_h - render_h) / 2
    else:
        render_h = box_h - 2 * margin
        render_w = render_h * img_ratio
        render_x = float(left) + (box_w - render_w) / 2
        render_y = float(top) + margin

    slide.shapes.add_picture(
        str(image_path),
        int(render_x),
        int(render_y),
        width=int(render_w),
        height=int(render_h),
    )


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)
    add_title_block(
        slide,
        "Sistema de Moeda Estudantil",
        "Arquitetura, tecnologias e fluxo completo de teste da aplicacao",
    )
    add_image_panel(slide, ASSETS_DIR / "11-showcase-focus.png", (7.35, 0.75, 5.45, 5.80))
    add_footer(slide, prs, "Equipe Lab03 | Release 1")


def add_card_box(slide, box: tuple[float, float, float, float], card: CardBlock) -> None:
    left, top, width, height = box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD_COLOR
    shape.line.color.rgb = MUTED_LINE
    shape.line.width = Pt(1.1)

    title_box = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(top + 0.18),
        Inches(width - 0.44),
        Inches(0.46),
    )
    title_tf = title_box.text_frame
    title_tf.clear()
    p_title = title_tf.paragraphs[0]
    p_title.text = card.title
    style_paragraph(p_title, size=17, color=TEXT_PRIMARY, bold=True)

    items_box = slide.shapes.add_textbox(
        Inches(left + 0.22),
        Inches(top + 0.66),
        Inches(width - 0.44),
        Inches(height - 0.82),
    )
    items_tf = items_box.text_frame
    items_tf.clear()
    items_tf.word_wrap = True

    for idx, item in enumerate(card.items):
        p = items_tf.paragraphs[0] if idx == 0 else items_tf.add_paragraph()
        p.text = f"- {item}"
        style_paragraph(p, size=14, color=TEXT_SECONDARY, line_spacing=1.12, space_after=4)


def add_cards_slide(
    prs: Presentation,
    *,
    title: str,
    subtitle: str,
    cards: list[CardBlock],
    index: int,
    total: int,
    footer: str = "Sistema de Moeda Estudantil | Lab03",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.55), Inches(9.0), Inches(0.65))
    title_tf = title_box.text_frame
    title_tf.clear()
    p_title = title_tf.paragraphs[0]
    p_title.text = title
    style_paragraph(p_title, size=31, color=TEXT_PRIMARY, bold=True)

    subtitle_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.16), Inches(10.8), Inches(0.45))
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.clear()
    p_sub = subtitle_tf.paragraphs[0]
    p_sub.text = subtitle
    style_paragraph(p_sub, size=15, color=TEXT_SECONDARY)

    boxes = [
        (0.62, 1.72, 6.10, 2.15),
        (6.95, 1.72, 5.75, 2.15),
        (0.62, 4.02, 6.10, 2.15),
        (6.95, 4.02, 5.75, 2.15),
    ]
    for card, box in zip(cards, boxes, strict=False):
        add_card_box(slide, box, card)

    add_footer(slide, prs, f"Slide {index}/{total} | {footer}")


def add_split_slide(prs: Presentation, slide_data: SplitSlide, index: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, prs)

    title_box = slide.shapes.add_textbox(Inches(0.62), Inches(0.55), Inches(7.2), Inches(0.68))
    title_tf = title_box.text_frame
    title_tf.clear()
    p_title = title_tf.paragraphs[0]
    p_title.text = slide_data.title
    style_paragraph(p_title, size=30, color=TEXT_PRIMARY, bold=True)

    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.62),
        Inches(1.45),
        Inches(7.00),
        Inches(5.35),
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = SURFACE_COLOR
    left_panel.line.color.rgb = MUTED_LINE
    left_panel.line.width = Pt(1.1)

    bullets_box = slide.shapes.add_textbox(Inches(0.95), Inches(1.80), Inches(6.35), Inches(4.80))
    bullets_tf = bullets_box.text_frame
    bullets_tf.clear()
    bullets_tf.word_wrap = True

    for idx, bullet in enumerate(slide_data.bullets):
        p = bullets_tf.paragraphs[0] if idx == 0 else bullets_tf.add_paragraph()
        p.text = f"• {bullet}"
        style_paragraph(p, size=18, color=TEXT_PRIMARY, line_spacing=1.18, space_after=8)

    if slide_data.image:
        image_path = ASSETS_DIR / slide_data.image
        box = IMG_BOXES.get(slide_data.image_mode, IMG_BOXES["landscape"])
        add_image_panel(slide, image_path, box)

    add_footer(slide, prs, f"Slide {index}/{total} | {slide_data.footer}")


def build_split_slides() -> list[SplitSlide]:
    return [
        SplitSlide(
            title="Problema e objetivo do produto",
            bullets=[
                "Reconhecimento academico era informal e sem trilha consolidada.",
                "Alunos nao tinham uma visao unica de saldo e historico.",
                "A proposta foi criar uma moeda digital com rastreabilidade.",
                "Resultado: incentivo real com resgate em empresas parceiras.",
            ],
            image="11-showcase-focus.png",
            image_mode="landscape",
        ),
        SplitSlide(
            title="Arquitetura ponta a ponta",
            bullets=[
                "Frontend SPA com experiencias separadas por papel.",
                "Backend em camadas com servicos e regras de negocio.",
                "Persistencia em MongoDB com ledger append-only.",
                "Comunicacao via API REST em /api com autenticacao por token.",
            ],
            image="07-diagrama-componentes.png",
            image_mode="diagram",
        ),
        SplitSlide(
            title="Fluxo 1: professor distribui moedas",
            bullets=[
                "Professor seleciona aluno, valor e mensagem de reconhecimento.",
                "Servico valida papel, saldo e consistencia do destino.",
                "Transacao e saldos sao atualizados no mesmo fluxo.",
                "Aluno recebe historico atualizado para acompanhamento.",
            ],
            image="14-prof-compositor-focus.png",
            image_mode="portrait",
        ),
        SplitSlide(
            title="Fluxo 2: aluno resgata com cupom",
            bullets=[
                "Aluno escolhe um beneficio ativo na vitrine.",
                "Sistema valida custo versus saldo disponivel.",
                "Cupom unico e emitido e registrado no ledger.",
                "Parceiro usa os dados para validar o atendimento.",
            ],
            image="13-aluno-resgate-focus.png",
            image_mode="portrait",
        ),
        SplitSlide(
            title="Operacao do parceiro",
            bullets=[
                "Parceiro atualiza perfil e publica novas vantagens.",
                "Catalogo permite ativar, desativar e excluir itens.",
                "Cada item define custo e descricao de forma objetiva.",
                "A vitrine publica passa a refletir o estado atual do board.",
            ],
            image="15-parceiro-board-focus.png",
            image_mode="portrait",
        ),
        SplitSlide(
            title="Persistencia e rastreabilidade",
            bullets=[
                "Entidades centrais: Student, Professor, Partner, Benefit.",
                "CoinTransaction registra ALLOCATION, TRANSFER e REDEMPTION.",
                "SessionToken controla autenticacao e expiracao de sessao.",
                "Allowance semestral garante idempotencia por professor.",
            ],
            image="09-modelo-er.png",
            image_mode="diagram",
        ),
        SplitSlide(
            title="Qualidade e evidencia de testes",
            bullets=[
                "Build backend e frontend validados antes da entrega.",
                "Teste E2E cobre autenticacao, autorizacao e regras criticas.",
                "Scripts Docker padronizam start, teste e limpeza completa.",
                "Demo reproduzivel em laboratorio com o mesmo passo a passo.",
            ],
            image="12-vitrine-focus.png",
            image_mode="landscape",
        ),
        SplitSlide(
            title="Conclusao e proximos passos",
            bullets=[
                "Release 1 entrega fluxo completo de ponta a ponta.",
                "Arquitetura modular facilita evolucao e manutencao.",
                "Proximos passos: SMTP real, observabilidade e hardening.",
                "Base pronta para feedback da banca e iteracoes do produto.",
            ],
            image="11-showcase-focus.png",
            image_mode="landscape",
            footer="Sistema de Moeda Estudantil | Obrigado",
        ),
    ]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate the Lab03 presentation deck.")
    parser.add_argument(
        "--output",
        type=Path,
        default=DEFAULT_OUTPUT,
        help="Output PPTX path (default: docs/slides/lab03s03-apresentacao.pptx)",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    split_slides = build_split_slides()
    total_slides = 1 + 3 + len(split_slides)

    add_title_slide(prs)

    add_cards_slide(
        prs,
        title="Tecnologias principais do frontend",
        subtitle="Foco em experiencia visual, tipagem e produtividade de entrega",
        cards=[
            CardBlock("Base da aplicacao", ["React 19", "TypeScript", "React Router"]),
            CardBlock("UI e interacao", ["Framer Motion", "Layout responsivo", "Componentizacao por papel"]),
            CardBlock("Consumo de API", ["Cliente HTTP tipado", "Fluxo de auth com token", "Tratamento de erros de rede"]),
            CardBlock("Build e qualidade", ["Vite", "ESLint", "Build de producao validado"]),
        ],
        index=2,
        total=total_slides,
    )

    add_cards_slide(
        prs,
        title="Tecnologias principais do backend",
        subtitle="Camadas bem separadas para regra de negocio, seguranca e persistencia",
        cards=[
            CardBlock("Core", ["Java 21", "Spring Boot 3.5", "Maven"]),
            CardBlock("Seguranca", ["Spring Security", "Token opaco", "RBAC via AuthFacade"]),
            CardBlock("Persistencia", ["MongoDB", "Spring Data", "Indices unicos e ledger"]),
            CardBlock("Robustez", ["GlobalExceptionHandler", "Servicos fail-fast", "Data seeder idempotente"]),
        ],
        index=3,
        total=total_slides,
    )

    add_cards_slide(
        prs,
        title="Integracao Docker para teste completo",
        subtitle="Pipeline unico para subir, validar rotas e limpar o ambiente",
        cards=[
            CardBlock("Subir stack", ["./docker/docker-start.sh", "Sobe frontend + backend + Mongo", "Ambiente padronizado para demo"]),
            CardBlock("Executar E2E", ["./docker/docker-e2e-routes-test.sh", "Cobre auth, role e regras", "Valida fluxos criticos"]),
            CardBlock("Encerrar", ["./docker/docker-stop-clean.sh", "Remove containers e volumes", "Limpa imagens locais"]),
            CardBlock("Beneficio", ["Execucao reproduzivel", "Mesmo roteiro para banca", "Feedback rapido por ciclo"]),
        ],
        index=4,
        total=total_slides,
    )

    for idx, data in enumerate(split_slides, start=5):
        add_split_slide(prs, data, idx, total_slides)

    output_path = args.output.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Presentation generated: {output_path}")
    print(f"Slides: {total_slides}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
